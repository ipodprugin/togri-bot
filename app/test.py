import pprint, jinja2, os, asyncio

from settings.config import settings
from pptx import Presentation  

pp = pprint.PrettyPrinter(indent=4)


models = [
    {
        'address': 'ЮВАО, район Печатники, проезд Проектируемый 5112, вл. 10, '
                   'этаж 1, помещение 3Н, комнаты 30-42, 44, 45',
        'applications_enddate': '15.01.2024',
        'auction_step': '80337.0',
        'deposit': '321348',
        'district_name': 'Не определен',
        'floor': 'Не указано',
        'procedure_form': 1,
        'm1_min_price': None,
        'm1_start_price': 40523.08,
        'min_price': 0.0,
        'object_area': 79.3,
        'price_decrease_step': '0.0',
        'region_name': 'ЮВАО',
        'start_price': '1606740.00',
        'subway_stations': None
    },
    {
        'address': 'улица Зои и Александра Космодемьянских, дом 4, корпус 3, '
                   'Подвал № 0',
        'applications_enddate': '15.01.2024',
        'auction_step': '271675.0',
        'deposit': '1086700',
        'district_name': 'Войковский район',
        'floor': 'Подвал',
        'procedure_form': 2,
        'm1_min_price': 33151.311775472845,
        'm1_start_price': 66302.62,
        'min_price': 5433500.0,
        'object_area': 163.9,
        'price_decrease_step': '543350.0',
        'region_name': 'САО',
        'start_price': '10867000.00',
        'subway_stations': None,
        'x': '5',
    }
]


async def render_text(input_path, model, output_path, jinja2_env):
    ppt = Presentation(input_path)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                await _render_text_frame(shape.text_frame, model, jinja2_env)
    ppt.save(output_path)


async def _render_text_frame(text_frame, model, jinja2_env):
    last_ok = True
    tmp_str = ""
    tmp_runs = []
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            cur_text = run.text
            tmp_str += cur_text
            tmp_runs.append(run)
            try:
                rtemplate = jinja2_env.from_string(tmp_str)
                rendered_text = rtemplate.render(model)

                for run in tmp_runs:
                    run.text = rendered_text
                    rendered_text = "" # overwrites text

                tmp_str = ""
                tmp_runs = []
                last_ok = True
            except Exception as e:
                #could not finish, i.e. have to append!
                last_ok = False


for index, model in enumerate(models):
    jinja2_env = jinja2.Environment()
    output_path = f'{settings.PPTX_OUTPUT_DIRPATH}/{index}.pptx'
    # asyncio.run(render_text('templates/template4.pptx', model, output_path, jinja2_env))
    asyncio.run(render_text(settings.PPTX_TEMPLATE_PATH, model, output_path, jinja2_env))
