import os
import jinja2

from ..data.models import SheetRowTenderContent

from settings.config import settings

from pptx import Presentation  
from collections import defaultdict


async def render_text(input_path, model, output_path, jinja2_env):
    ppt = Presentation(input_path)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                await _render_text_frame(shape.text_frame, model, jinja2_env)
    ppt.save(output_path)


async def _render_text_frame(text_frame, model, jinja2_env):
    last_ok = True
    tmp_str = ""
    tmp_runs = []
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            cur_text = run.text
            tmp_str += cur_text
            tmp_runs.append(run)
            try:
                rtemplate = jinja2_env.from_string(tmp_str)
                rendered_text = rtemplate.render(model)

                for run in tmp_runs:
                    run.text = rendered_text
                    rendered_text = "" # overwrites text

                tmp_str = ""
                tmp_runs = []
                last_ok = True
            except Exception as e:
                #could not finish, i.e. have to append!
                last_ok = False


def delete_slide(presentation,  index):
    xml_slides = presentation.slides._sldIdLst  # pylint: disable=W0212
    slides = list(xml_slides)
    xml_slides.remove(slides[index])


async def replace_images_by_shape_text(images: dict, template_path: str, output_path: str, images_slides_ids: list[int] = [2, 3, 4]):
    # TODO: 
    #   - ОТРЕФАКТОРИТЬ ЭТОТ УЖАС
    #   - Удалять плейсхолдер shape после вставки картинки
    #   - Сделать фиксированный словарь плейсхолдеров и итерироваться по нему
    #   - Удалять плейсхолдер, если для него нет картинки
    prs = Presentation(template_path)
    found_index = -1
    images_slides_ids_dict = defaultdict(int)
    for image in images:
        image_file = images[image]
        search_str = image
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    found_index = shape.text.find(search_str)
                    if found_index != -1:
                        if image_file:
                            slide_id = prs.slides.index(slide)
                            if slide_id in images_slides_ids:
                                images_slides_ids_dict[slide_id] += 1
                            horiz_ = shape.left
                            vert_ = shape.top
                            height_ = shape.height
                            width_ = shape.width
                            slide.shapes.add_picture(image_file, horiz_, vert_, width_, height_)
                        sp = shape._sp
                        sp.getparent().remove(sp)
                        break
                if found_index != -1:
                    break
    
    for slide_id in images_slides_ids_dict:
        if images_slides_ids_dict[slide_id] < 3:
            delete_slide(prs, slide_id)

    prs.save(output_path)


async def render_pptx(tender: SheetRowTenderContent, pictures: dict):
    print('------------ tender', tender)
    model = {
        "address": tender.address,
        "subway_stations": tender.subway_stations,
        "region_name": tender.region_name,
        "district_name": tender.district_name,
        "object_area": tender.object_area,
        "floor": tender.floor,
        "applications_enddate": tender.applications_enddate,
        "deposit": tender.deposit,
        "start_price": tender.start_price,
        "m1_start_price": tender.m1_start_price,
        "min_price": tender.min_price,
        "m1_min_price": tender.m1_min_price,
        "procedure_form": tender.procedure_form,
        "auction_step": tender.auction_step,
        "price_decrease_step": tender.price_decrease_step,
    }

    print('------------ model', model)

    jinja2_env = jinja2.Environment()

    if not os.path.isdir(settings.PPTX_OUTPUT_DIRPATH):
        os.mkdir(settings.PPTX_OUTPUT_DIRPATH)
    output_path = f'{settings.PPTX_OUTPUT_DIRPATH}/{tender.id}.pptx'
    await render_text(settings.PPTX_TEMPLATE_PATH, model, output_path, jinja2_env)
    await replace_images_by_shape_text(images=pictures, template_path=output_path, output_path=output_path)
    return output_path
